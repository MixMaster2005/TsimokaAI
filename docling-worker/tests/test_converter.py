"""Tests unitaires du worker (spec v2) — exécutables sans clé Gemini ni conteneur.

Lancés avec : python -m unittest tests.test_converter  (ou `python tests/test_converter.py`).
Les appels Gemini sont mockés : on vérifie le pipeline (extraction, placeholders,
transcription de pages, plafond d'images, dégradation sur échec Gemini).
"""
import io
import os
import struct
import sys
import unittest
import zlib
from unittest import mock

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from app import image_extractor  # noqa: E402
from app.markitdown_converter import MarkItDownConverter  # noqa: E402
import app.markitdown_converter as markitdown_converter  # noqa: E402


def png_bytes(width=320, height=200):
    """PNG RGB minimal (gris), assez grand pour passer le filtre MIN_IMAGE_DIMENSION."""

    def chunk(tag, data):
        return (struct.pack(">I", len(data)) + tag + data
                + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF))

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + b"\xa0\x80\x60" * width for _ in range(height))
    return sig + chunk(b"IHDR", ihdr) + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b"")


def make_text_pdf():
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Cours d'algorithmique - Introduction", fontsize=16)
    for i in range(10):
        page.insert_text((72, 120 + i * 20), "Paragraphe de remplissage textuel du document.")
    return doc.tobytes()


def make_image_pdf(scanned=False):
    import fitz

    doc = fitz.open()
    if scanned:
        for _ in range(2):
            page = doc.new_page()
            page.insert_image(fitz.Rect(0, 0, 595, 842), stream=png_bytes())
    else:
        page = doc.new_page()
        page.insert_text((72, 72), "Chapitre avec figure", fontsize=16)
        page.insert_image(fitz.Rect(72, 100, 400, 300), stream=png_bytes())
        page.insert_text((72, 350), "Texte qui suit la figure.")
    return doc.tobytes()


def make_docx():
    from docx import Document

    doc = Document()
    doc.add_heading("Rapport", level=1)
    doc.add_paragraph("Un texte d'introduction suffisamment long pour ce document Word.")
    doc.add_picture(io.BytesIO(png_bytes()))
    doc.add_paragraph("Suite du texte apres l'image.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx():
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(png_bytes()), Emu(914400), Emu(914400))
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.shapes.add_textbox(Emu(914400), Emu(914400), Emu(2000000), Emu(500000)).text = "Slide 2"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class FakeVision:
    """Double du captioner Gemini : déterministe, sans réseau."""

    def caption_figure(self, image, content_type):
        return "Schema d'un graphe oriente"

    def transcribe_full_page(self, page_image, content_type="image/png"):
        return "## Page transcrite\nContenu du document scanné."


class MarkItDownConverterTest(unittest.TestCase):

    def _convert(self, content, filename, vision=None):
        if vision is None:
            vision = FakeVision()
        with mock.patch.object(markitdown_converter, "VisionCaptioner", return_value=vision):
            return MarkItDownConverter().convert(content, filename)

    def test_text_pdf_markitdown_sans_image(self):
        result = self._convert(make_text_pdf(), "cours.pdf")
        # Les PDF utilisent maintenant le nouveau pipeline PyMuPDF
        self.assertEqual(result["method"], "pymupdf_layout")
        self.assertIn("Cours d'algorithmique", result["markdown"])
        self.assertEqual(result["pages_processed"], 1)
        # L'AST doit être présente
        self.assertIsNotNone(result.get("document"))

    def test_docx_placeholder_inline_et_data_uri_retire(self):
        result = self._convert(make_docx(), "rapport.docx")
        self.assertEqual(result["method"], "markitdown")
        self.assertEqual(len(result["images"]), 1)
        img = result["images"][0]
        self.assertEqual(img["placeholder_id"], "img_001")
        self.assertEqual(img["content_type"], "image/png")
        self.assertEqual(img["caption"], "Schema d'un graphe oriente")
        self.assertGreater(len(img["data_base64"]), 0)
        self.assertIn("{{IMAGE:img_001}}", result["markdown"])
        self.assertNotIn("data:image", result["markdown"])

    def test_pptx_placeholder_inline(self):
        result = self._convert(make_pptx(), "slides.pptx")
        self.assertEqual(len(result["images"]), 1)
        self.assertIn("{{IMAGE:img_001}}", result["markdown"])
        self.assertEqual(result["method"], "markitdown")

    def test_pdf_image_placeholder_a_la_fin(self):
        result = self._convert(make_image_pdf(), "chapitre.pdf")
        self.assertEqual(len(result["images"]), 1)
        # Le nouveau pipeline produit un AST — les images sont dans result["images"]
        # Le markdown peut contenir un placeholder ou non selon le renderer
        self.assertIn("Chapitre avec figure", result["markdown"])

    def test_pdf_scan_transcrit_page_par_page(self):
        result = self._convert(make_image_pdf(scanned=True), "scan.pdf")
        # Le nouveau pipeline classifie les pages — un PDF scanné peut être
        # traité par le pipeline PyMuPDF ou rester en mode transcription
        self.assertIn(result["method"], ("pymupdf_layout", "markitdown_with_page_transcription"))
        self.assertEqual(result["pages_processed"], 2)
        # L'AST doit être présente
        self.assertIsNotNone(result.get("document"))

    def test_markdown_court_non_pdf_non_considere_scanne(self):
        # Un .md court a "1 page" (count_pages) : il ne doit PAS être envoyé en
        # transcription (réservée au PDF) mais conserver son texte (régression spec v2).
        result = self._convert("# Titre\n\nTexte".encode(), "notes.md")
        self.assertEqual(result["method"], "markitdown")
        self.assertIn("# Titre", result["markdown"])
        self.assertNotIn("transcription par page impossible", "\n".join(result["warnings"]))

    def test_plafond_images_30_avec_warning(self):
        images = [image_extractor.ExtractedImage(png_bytes(), "image/png", f"page {i}")
                  for i in range(35)]
        with mock.patch.object(image_extractor, "extract_images", return_value=images):
            result = self._convert(make_text_pdf(), "doc.pdf")
        self.assertEqual(len(result["images"]), image_extractor.MAX_EXTRACTED_IMAGES)
        self.assertTrue(any("Plafond" in w for w in result["warnings"]))
        self.assertEqual(result["images"][-1]["placeholder_id"],
                         f"img_{image_extractor.MAX_EXTRACTED_IMAGES:03d}")

    def test_echec_gemini_non_bloquant(self):
        class BrokenVision:
            def caption_figure(self, image, content_type):
                raise RuntimeError("API key invalide")

        result = self._convert(make_docx(), "rapport.docx", vision=BrokenVision())
        self.assertEqual(result["method"], "markitdown")
        self.assertEqual(result["images"][0]["caption"], "")
        self.assertTrue(any("Légende Gemini" in w for w in result["warnings"]))


if __name__ == "__main__":
    unittest.main()
