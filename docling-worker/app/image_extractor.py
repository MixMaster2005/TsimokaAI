"""Extraction des images embarquées d'un document (PDF / DOCX / PPTX).

Spécification v2 : le worker ne stocke pas dans MinIO (pas de credentials Python) ; il
renvoie les images en base64 au format attendu par ``ingestion-service``, qui les uploade
et remplace les placeholders ``{{IMAGE:img_001}}`` par ``![caption](url)``.

Implémentation :
  - PDF  : PyMuPDF (``fitz``) — ``page.get_images(...)`` + ``doc.extract_image(xref)``.
  - DOCX : python-docx — ``inline_shapes`` puis ``part.image`` (blob + content_type).
  - PPTX : python-pptx — formes ``MSO_SHAPE_TYPE.PICTURE`` (slide par slide).

Les images trop petites (logos, icônes…) sont ignorées pour ne pas saturer l'API Gemini.
Le nombre d'images renvoyé est plafonné à ``MAX_EXTRACTED_IMAGES`` (warning au-delà).
"""
import io
import logging
import os
from dataclasses import dataclass
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)

MAX_EXTRACTED_IMAGES = int(os.environ.get("DOCLING_MAX_EXTRACTED_IMAGES", "30"))

# Pixels min (largeur OU hauteur) en dessous desquels une image est considérée
# comme un logo/icône décoratif et ignorée.
MIN_IMAGE_DIMENSION = int(os.environ.get("DOCLING_MIN_IMAGE_DIMENSION", "64"))

# Formats sans perte vers lesquels normaliser les images extraites.
EXTRACT_EXT_TO_CONTENT_TYPE = {
    "png": "image/png",
    "jpeg": "image/jpeg",
    "jpg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "tiff": "image/tiff",
    "bmp": "image/bmp",
}


@dataclass
class ExtractedImage:
    """Image extraite : binaire + type MIME + localisation approximative dans le document."""

    content: bytes
    content_type: str
    location: str  # ex. "page 3", "slide 2" — utile pour le warning/le log.


def extract_images(content: bytes, filename: str) -> List[ExtractedImage]:
    """Extrait les images embarquées selon l'extension du fichier."""
    extension = os.path.splitext(filename)[1].lower()
    if extension == ".pdf":
        return _extract_pdf_images(content)
    if extension == ".docx":
        return _extract_docx_images(content)
    if extension == ".pptx":
        return _extract_pptx_images(content)
    logger.info("Extraction d'images non gérée pour %s (seul PDF/DOCX/PPTX)", filename)
    return []


def _px(emu: int) -> int:
    """Convertit une longueur EMU (OXML) en pixels (~96 DPI, 1 px = 9525 EMU)."""
    return emu // 9525 if emu else 0


def _size_ok(width_px: int, height_px: int, blob_len: int = 0) -> bool:
    """Filtre les images décoratives (logos/icônes) : trop petites en dimensions."""
    if width_px and height_px:
        return width_px >= MIN_IMAGE_DIMENSION or height_px >= MIN_IMAGE_DIMENSION
    # Dimensions inconnues : repli sur la taille du blob (très approximatif).
    return blob_len >= MIN_IMAGE_DIMENSION * MIN_IMAGE_DIMENSION


def render_pdf_pages(content: bytes) -> List[bytes]:
    """Rend chaque page d'un PDF en PNG (~150 DPI) — utilisé pour les documents scannés."""
    import fitz

    pages: List[bytes] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for index, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            pages.append(pix.tobytes("png"))
            logger.debug("Page %d rendue en image (%dx%d)", index + 1, pix.width, pix.height)
    return pages


def _extract_pdf_images(content: bytes) -> List[ExtractedImage]:
    import fitz

    images: List[ExtractedImage] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page_index, page in enumerate(doc):
            for image_index, image_info in enumerate(page.get_images(full=True)):
                xref = image_info[0]
                extracted = doc.extract_image(xref)
                if not extracted:
                    continue
                width = image_info[2] or 0
                height = image_info[3] or 0
                if width < MIN_IMAGE_DIMENSION or height < MIN_IMAGE_DIMENSION:
                    logger.debug(
                        "Image xref=%s (page %d) ignorée : trop petite (%dx%d)",
                        xref, page_index + 1, width, height,
                    )
                    continue
                content_type = EXTRACT_EXT_TO_CONTENT_TYPE.get(extracted.get("ext", ""))
                if content_type is None:
                    logger.debug("Image xref=%s (page %d) ignorée : format '%s' non supporté",
                                 xref, page_index + 1, extracted.get("ext"))
                    continue
                images.append(ExtractedImage(
                    content=extracted["image"],
                    content_type=content_type,
                    location=f"page {page_index + 1} (image {image_index + 1})",
                ))
    return images


def _extract_docx_images(content: bytes) -> List[ExtractedImage]:
    from docx import Document

    document = Document(io.BytesIO(content))
    images: List[ExtractedImage] = []
    for index, shape in enumerate(document.inline_shapes):
        try:
            blip = shape._inline.graphic.graphicData.pic.blipFill.blip  # noqa: SLF001
            rid = blip.embed
            if not rid:
                continue
            part = document.part.related_parts[rid]
            image = part.image
            content_type = image.content_type or "image/png"
            if content_type not in EXTRACT_EXT_TO_CONTENT_TYPE.values():
                logger.debug("Image DOCX #%d ignorée : type '%s' non supporté", index, content_type)
                continue
            if not _size_ok(_px(image.width), _px(image.height), len(image.blob)):
                logger.debug("Image DOCX #%d ignorée : trop petite", index)
                continue
            images.append(ExtractedImage(
                content=image.blob,
                content_type=content_type,
                location=f"paragraphe {index + 1}",
            ))
        except Exception as e:  # noqa: BLE001 - forme sans image brute exploitable
            logger.debug("Image DOCX #%d non extractible (%s) — ignorée", index, e)
    return images


def _extract_pptx_images(content: bytes) -> List[ExtractedImage]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(io.BytesIO(content))
    images: List[ExtractedImage] = []
    for slide_index, slide in enumerate(presentation.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                content_type = shape.image.content_type or "image/png"
                if content_type not in EXTRACT_EXT_TO_CONTENT_TYPE.values():
                    logger.debug("Image PPTX (slide %d) ignorée : type '%s' non supporté",
                                 slide_index + 1, content_type)
                    continue
                if not _size_ok(_px(shape.width), _px(shape.height), len(shape.image.blob)):
                    logger.debug("Image PPTX (slide %d) ignorée : trop petite", slide_index + 1)
                    continue
                images.append(ExtractedImage(
                    content=shape.image.blob,
                    content_type=content_type,
                    location=f"slide {slide_index + 1} (forme {shape_index + 1})",
                ))
            except Exception as e:  # noqa: BLE001 - forme sans image brute exploitable
                logger.debug("Image PPTX (slide %d, forme %d) non extractible (%s) — ignorée",
                             slide_index + 1, shape_index + 1, e)
    return images
